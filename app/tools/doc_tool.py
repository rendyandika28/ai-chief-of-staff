"""Generate dokumen (md/docx/pptx) dari markdown, balikin sebagai file attachment.

Konten markdown dibuat oleh LLM; tool ini cuma render ke format yang diminta,
simpen ke temp file, balikin marker [FILE:path] yang ditangkep telegram.py.
"""

import os
import re
import tempfile
import time


def _slug(s: str) -> str:
    s = re.sub(r"[^\w\s-]", "", s).strip().lower()
    return re.sub(r"[\s_-]+", "-", s)[:60] or "dokumen"


def _blocks(content: str):
    """Pecah markdown jadi list (heading_slide, [baris_body]) per H1 '# '."""
    blocks, cur = [], None
    for line in content.splitlines():
        s = line.strip()
        if s.startswith("# "):
            if cur:
                blocks.append(cur)
            cur = (s[2:].strip(), [])
        elif s:
            if cur is None:
                cur = ("", [])
            for p in ("### ", "## ", "- ", "* "):
                if s.startswith(p):
                    s = s[len(p):]
                    break
            cur[1].append(s)
    if cur:
        blocks.append(cur)
    return blocks


def _write_md(content: str, path: str):
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)


def _write_docx(content: str, title: str, path: str):
    from docx import Document
    doc = Document()
    doc.add_heading(title, level=0)
    for line in content.splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith("### "):
            doc.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s.startswith(("- ", "* ")):
            doc.add_paragraph(s[2:], style="List Bullet")
        else:
            doc.add_paragraph(s)
    doc.save(path)


def _write_pptx(content: str, title: str, path: str):
    from pptx import Presentation
    prs = Presentation()
    blocks = _blocks(content) or [(title, [])]
    layout = prs.slide_layouts[1]  # Title and Content
    for stitle, lines in blocks:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = stitle or title
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, ln in enumerate(lines):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = ln
    prs.save(path)


_WRITERS = {"md": _write_md, "docx": _write_docx, "pptx": _write_pptx}


class DocTool:
    name = "doc_gen"
    description = (
        "Bikin dokumen file (brief, kontrak, laporan, atau slide) yang dikirim ke user "
        "sebagai attachment di Telegram. Lo yang nulis isinya dalam markdown.\n"
        "Format input: <format>|<judul>|<isi markdown>\n"
        "  - format: md | docx | pptx\n"
        "  - judul: judul dokumen / nama file (tanpa ekstensi)\n"
        "  - isi: markdown. Pake '# ' buat heading (di pptx tiap '# ' = 1 slide baru), "
        "'## '/'### ' subheading, '- ' buat bullet.\n"
        "Contoh: docx|Brief Project X|# Problem\\n- poin a\\n- poin b\\n# Ide Inti\\n...\n"
        "Pilih docx buat kontrak/brief formal, pptx buat presentasi, md buat catatan cepet."
    )

    def run(self, input: str = "") -> str:
        parts = (input or "").split("|", 2)
        if len(parts) != 3:
            return "Error: format harus <md|docx|pptx>|<judul>|<isi markdown>"
        fmt, title, content = parts[0].strip().lower(), parts[1].strip(), parts[2]
        if fmt not in _WRITERS:
            return f"Error: format '{fmt}' gak didukung (pilih md/docx/pptx)"
        if not content.strip():
            return "Error: isi dokumen kosong"

        fname = f"{_slug(title)}-{int(time.time())}.{fmt}"
        path = os.path.join(tempfile.gettempdir(), fname)
        try:
            if fmt == "md":
                _write_md(content, path)
            elif fmt == "docx":
                _write_docx(content, title or "Dokumen", path)
            else:
                _write_pptx(content, title or "Presentasi", path)
        except Exception as e:
            return f"Error bikin dokumen: {e}"

        return f"Dokumen '{title}' ({fmt}) udah jadi. [FILE:{path}]"


def _demo():
    t = DocTool()
    md = "# Problem\n- users bingung onboarding\n# Ide Inti\nAsisten AI yang nuntun step-by-step."
    for fmt in ("md", "docx", "pptx"):
        out = t.run(f"{fmt}|Brief Test|{md}")
        m = re.search(r"\[FILE:(.*?)\]", out)
        assert m, f"{fmt}: no marker in {out!r}"
        assert os.path.exists(m.group(1)), f"{fmt}: file not created"
        os.remove(m.group(1))
    # input jelek
    assert t.run("cuma-satu-bagian").startswith("Error")
    assert t.run("xls|Judul|isi").startswith("Error")
    print("doc_tool OK")


if __name__ == "__main__":
    _demo()
